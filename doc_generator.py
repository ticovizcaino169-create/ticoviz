"""
TicoViz Corporation v2 — Generador de Documentos
Crea PDFs profesionales, presentaciones PPTX y hojas Excel.
"""
import logging
import os
from pathlib import Path
from datetime import datetime

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    PageBreak, HRFlowable
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

from config import PRODUCTS_DIR

logger = logging.getLogger("ticoviz.docs")

# Colores corporativos TicoViz
COLORS = {
    "primary": "#1a1a2e",
    "secondary": "#16213e",
    "accent": "#0f3460",
    "highlight": "#e94560",
    "text": "#333333",
    "light": "#f0f0f0",
    "white": "#ffffff",
}


def _ensure_product_dir(product_id: int) -> Path:
    product_dir = PRODUCTS_DIR / f"product_{product_id}"
    product_dir.mkdir(parents=True, exist_ok=True)
    return product_dir


# ==================== PDF ====================

def generar_pdf(product_id: int, datos: dict) -> str:
    """Genera un PDF profesional con documentacion completa del producto."""
    product_dir = _ensure_product_dir(product_id)
    safe_name = datos.get("nombre", "producto").replace(" ", "_").replace("/", "_")[:50]
    filepath = str(product_dir / f"{safe_name}.pdf")

    doc = SimpleDocTemplate(
        filepath, pagesize=letter,
        rightMargin=72, leftMargin=72,
        topMargin=72, bottomMargin=72,
    )

    styles = getSampleStyleSheet()

    styles.add(ParagraphStyle(
        name="TicoTitle", parent=styles["Title"],
        fontSize=28, textColor=HexColor(COLORS["primary"]),
        spaceAfter=20, alignment=TA_CENTER,
    ))
    styles.add(ParagraphStyle(
        name="TicoSubtitle", parent=styles["Normal"],
        fontSize=14, textColor=HexColor(COLORS["accent"]),
        spaceAfter=30, alignment=TA_CENTER,
    ))
    styles.add(ParagraphStyle(
        name="TicoHeading", parent=styles["Heading1"],
        fontSize=18, textColor=HexColor(COLORS["primary"]),
        spaceBefore=20, spaceAfter=10,
    ))
    styles.add(ParagraphStyle(
        name="TicoBody", parent=styles["Normal"],
        fontSize=11, textColor=HexColor(COLORS["text"]),
        leading=16, alignment=TA_JUSTIFY, spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        name="TicoHighlight", parent=styles["Normal"],
        fontSize=12, textColor=HexColor(COLORS["highlight"]),
        spaceBefore=5, spaceAfter=5,
    ))

    elements = []

    # Portada
    elements.append(Spacer(1, 2 * inch))
    elements.append(Paragraph("TicoViz Corporation", styles["TicoSubtitle"]))
    elements.append(Paragraph(datos.get("nombre", "Producto Digital"), styles["TicoTitle"]))
    elements.append(Spacer(1, 0.3 * inch))
    elements.append(Paragraph(
        datos.get("descripcion", ""),
        styles["TicoBody"],
    ))
    elements.append(Spacer(1, 0.5 * inch))
    elements.append(Paragraph(
        f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}",
        ParagraphStyle("Date", parent=styles["Normal"],
                       fontSize=10, textColor=HexColor("#888888"),
                       alignment=TA_CENTER),
    ))
    elements.append(PageBreak())

    # Seccion: Resumen Ejecutivo
    elements.append(Paragraph("Resumen Ejecutivo", styles["TicoHeading"]))
    elements.append(HRFlowable(
        width="100%", thickness=2,
        color=HexColor(COLORS["highlight"]),
        spaceAfter=10,
    ))
    desc_larga = datos.get("descripcion_larga", datos.get("descripcion", ""))
    for parrafo in desc_larga.split("\n"):
        if parrafo.strip():
            elements.append(Paragraph(parrafo.strip(), styles["TicoBody"]))

    # Seccion: Especificaciones Tecnicas
    elements.append(Spacer(1, 0.3 * inch))
    elements.append(Paragraph("Especificaciones Tecnicas", styles["TicoHeading"]))
    elements.append(HRFlowable(
        width="100%", thickness=2,
        color=HexColor(COLORS["highlight"]),
        spaceAfter=10,
    ))

    specs_data = [
        ["Propiedad", "Valor"],
        ["Categoria", datos.get("categoria", "N/A")],
        ["Tecnologias", datos.get("tecnologias", "N/A")],
        ["Publico Objetivo", datos.get("publico_objetivo", "N/A")],
        ["Potencial", datos.get("potencial", "N/A")],
        ["Precio Sugerido", f"${datos.get('precio_sugerido', 0):.2f} USD"],
    ]

    specs_table = Table(specs_data, colWidths=[2.5 * inch, 4 * inch])
    specs_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), HexColor(COLORS["primary"])),
        ("TEXTCOLOR", (0, 0), (-1, 0), HexColor(COLORS["white"])),
        ("FONTSIZE", (0, 0), (-1, 0), 12),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
        ("BACKGROUND", (0, 1), (-1, -1), HexColor("#f8f9fa")),
        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#dee2e6")),
        ("PADDING", (0, 0), (-1, -1), 8),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    elements.append(specs_table)

    # Seccion: Features
    features = datos.get("features_clave", [])
    if features:
        elements.append(Spacer(1, 0.3 * inch))
        elements.append(Paragraph("Funcionalidades Clave", styles["TicoHeading"]))
        elements.append(HRFlowable(
            width="100%", thickness=2,
            color=HexColor(COLORS["highlight"]),
            spaceAfter=10,
        ))
        for i, feat in enumerate(features, 1):
            elements.append(Paragraph(
                f"<b>{i}.</b> {feat}",
                styles["TicoBody"],
            ))

    # Seccion: Competencia
    competencia = datos.get("competencia", "")
    if competencia:
        elements.append(Spacer(1, 0.3 * inch))
        elements.append(Paragraph("Analisis de Competencia", styles["TicoHeading"]))
        elements.append(HRFlowable(
            width="100%", thickness=2,
            color=HexColor(COLORS["highlight"]),
            spaceAfter=10,
        ))
        elements.append(Paragraph(competencia, styles["TicoBody"]))

    # Footer
    elements.append(Spacer(1, 0.5 * inch))
    elements.append(HRFlowable(width="100%", thickness=1, color=HexColor("#cccccc")))
    elements.append(Paragraph(
        "TicoViz Corporation — Productos Digitales con IA",
        ParagraphStyle("Footer", parent=styles["Normal"],
                       fontSize=9, textColor=HexColor("#999999"),
                       alignment=TA_CENTER, spaceBefore=10),
    ))

    doc.build(elements)
    logger.info(f"PDF generado: {filepath}")
    return filepath


# ==================== PPTX ====================

def generar_pptx(product_id: int, datos: dict) -> str:
    """Genera una presentacion profesional PPTX."""
    product_dir = _ensure_product_dir(product_id)
    safe_name = datos.get("nombre", "producto").replace(" ", "_").replace("/", "_")[:50]
    filepath = str(product_dir / f"{safe_name}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary_rgb = RGBColor(0x1a, 0x1a, 0x2e)
    accent_rgb = RGBColor(0x0f, 0x34, 0x60)
    highlight_rgb = RGBColor(0xe9, 0x45, 0x60)
    white_rgb = RGBColor(0xff, 0xff, 0xff)
    light_rgb = RGBColor(0xf0, 0xf0, 0xf0)

    def add_slide(layout_idx=6):
        """Agrega slide con fondo oscuro."""
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = primary_rgb
        return slide

    def add_text_box(slide, left, top, width, height, text,
                     font_size=18, color=white_rgb, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return txBox

    # Slide 1: Portada
    slide = add_slide()
    add_text_box(slide, 1, 0.5, 11, 1, "TicoViz Corporation",
                 font_size=16, color=highlight_rgb, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.5, 11, 2, datos.get("nombre", "Producto Digital"),
                 font_size=44, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 4.8, 9, 1.5, datos.get("descripcion", ""),
                 font_size=18, color=light_rgb, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.5, 11, 0.5,
                 f"Generado: {datetime.now().strftime('%d/%m/%Y')}",
                 font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

    # Slide 2: Resumen
    slide = add_slide()
    add_text_box(slide, 0.8, 0.5, 11, 1, "Resumen del Producto",
                 font_size=32, bold=True, color=highlight_rgb)
    desc = datos.get("descripcion_larga", datos.get("descripcion", ""))
    add_text_box(slide, 0.8, 1.8, 11, 5, desc[:600],
                 font_size=16, color=light_rgb)

    # Slide 3: Especificaciones
    slide = add_slide()
    add_text_box(slide, 0.8, 0.5, 11, 1, "Especificaciones Tecnicas",
                 font_size=32, bold=True, color=highlight_rgb)
    specs_text = (
        f"Categoria: {datos.get('categoria', 'N/A')}\n\n"
        f"Tecnologias: {datos.get('tecnologias', 'N/A')}\n\n"
        f"Publico: {datos.get('publico_objetivo', 'N/A')}\n\n"
        f"Potencial: {datos.get('potencial', 'N/A')}\n\n"
        f"Precio: ${datos.get('precio_sugerido', 0):.2f} USD"
    )
    add_text_box(slide, 0.8, 1.8, 11, 5, specs_text,
                 font_size=20, color=light_rgb)

    # Slide 4: Features
    features = datos.get("features_clave", [])
    if features:
        slide = add_slide()
        add_text_box(slide, 0.8, 0.5, 11, 1, "Funcionalidades Clave",
                     font_size=32, bold=True, color=highlight_rgb)
        feat_text = "\n\n".join([f"  {i+1}. {f}" for i, f in enumerate(features)])
        add_text_box(slide, 0.8, 1.8, 11, 5, feat_text,
                     font_size=18, color=light_rgb)

    # Slide 5: Contacto
    slide = add_slide()
    add_text_box(slide, 1, 2.5, 11, 2, "Gracias",
                 font_size=48, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 11, 1, "TicoViz Corporation — Productos Digitales con IA",
                 font_size=18, color=highlight_rgb, alignment=PP_ALIGN.CENTER)

    prs.save(filepath)
    logger.info(f"PPTX generado: {filepath}")
    return filepath


# ==================== EXCEL ====================

def generar_excel_financiero(product_id: int, datos: dict) -> str:
    """Genera hoja Excel con analisis financiero del producto."""
    product_dir = _ensure_product_dir(product_id)
    safe_name = datos.get("nombre", "producto").replace(" ", "_").replace("/", "_")[:50]
    filepath = str(product_dir / f"{safe_name}_financiero.xlsx")

    wb = Workbook()

    # Estilos
    header_font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1A1A2E", end_color="1A1A2E", fill_type="solid")
    accent_fill = PatternFill(start_color="E94560", end_color="E94560", fill_type="solid")
    light_fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
    border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )
    center_align = Alignment(horizontal="center", vertical="center")

    # Hoja 1: Resumen
    ws = wb.active
    ws.title = "Resumen"
    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 45

    rows = [
        ("Propiedad", "Valor"),
        ("Producto", datos.get("nombre", "")),
        ("Categoria", datos.get("categoria", "")),
        ("Tecnologias", datos.get("tecnologias", "")),
        ("Precio Sugerido (USD)", datos.get("precio_sugerido", 0)),
        ("Potencial", datos.get("potencial", "")),
        ("Publico Objetivo", datos.get("publico_objetivo", "")),
        ("Fecha", datetime.now().strftime("%d/%m/%Y %H:%M")),
    ]

    for row_idx, (key, val) in enumerate(rows, 1):
        ws.cell(row=row_idx, column=1, value=key)
        ws.cell(row=row_idx, column=2, value=val)
        if row_idx == 1:
            for col in [1, 2]:
                cell = ws.cell(row=row_idx, column=col)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = center_align
        else:
            ws.cell(row=row_idx, column=1).font = Font(bold=True)
        for col in [1, 2]:
            ws.cell(row=row_idx, column=col).border = border

    # Hoja 2: Proyeccion financiera
    ws2 = wb.create_sheet("Proyeccion")
    ws2.column_dimensions["A"].width = 15
    for c in ["B", "C", "D", "E"]:
        ws2.column_dimensions[c].width = 18

    precio = float(datos.get("precio_sugerido", 50))
    headers = ["Mes", "Ventas Est.", "Ingreso", "Costos", "Ganancia"]
    for col, h in enumerate(headers, 1):
        cell = ws2.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = border

    for mes in range(1, 13):
        ventas = mes * 2 + 3  # Proyeccion simple
        ingreso = ventas * precio
        costos = ingreso * 0.15  # 15% costos
        ganancia = ingreso - costos

        row = mes + 1
        values = [f"Mes {mes}", ventas, ingreso, costos, ganancia]
        for col, val in enumerate(values, 1):
            cell = ws2.cell(row=row, column=col, value=val)
            cell.border = border
            if col >= 3:
                cell.number_format = "$#,##0.00"
            if mes % 2 == 0:
                cell.fill = light_fill

    # Totales
    total_row = 14
    ws2.cell(row=total_row, column=1, value="TOTAL").font = Font(bold=True, size=12)
    for col in range(2, 6):
        cell = ws2.cell(row=total_row, column=col)
        col_letter = chr(64 + col)
        cell.value = f"=SUM({col_letter}2:{col_letter}13)"
        cell.font = Font(bold=True, size=12)
        cell.fill = accent_fill
        cell.font = Font(bold=True, size=12, color="FFFFFF")
        cell.border = border
        if col >= 3:
            cell.number_format = "$#,##0.00"

    wb.save(filepath)
    logger.info(f"Excel generado: {filepath}")
    return filepath


# ==================== PAQUETE COMPLETO ====================

def generar_paquete_completo(product_id: int, datos: dict) -> dict:
    """
    Genera todos los documentos para un producto:
    - PDF con documentacion
    - PPTX con presentacion
    - Excel con analisis financiero
    """
    rutas = {}
    try:
        rutas["pdf"] = generar_pdf(product_id, datos)
        logger.info(f"PDF OK para producto #{product_id}")
    except Exception as e:
        logger.error(f"Error generando PDF #{product_id}: {e}")
        rutas["pdf"] = ""

    try:
        rutas["pptx"] = generar_pptx(product_id, datos)
        logger.info(f"PPTX OK para producto #{product_id}")
    except Exception as e:
        logger.error(f"Error generando PPTX #{product_id}: {e}")
        rutas["pptx"] = ""

    try:
        rutas["xlsx"] = generar_excel_financiero(product_id, datos)
        logger.info(f"Excel OK para producto #{product_id}")
    except Exception as e:
        logger.error(f"Error generando Excel #{product_id}: {e}")
        rutas["xlsx"] = ""

    return rutas
